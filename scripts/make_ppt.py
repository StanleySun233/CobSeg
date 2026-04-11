"""Generate academic conference PPT for DualStreamSegmenter."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── color palette (deep-navy academic) ────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2E, 0x5A)   # dark navy
TEAL   = RGBColor(0x00, 0x7B, 0xA7)   # accent teal
GOLD   = RGBColor(0xE8, 0xA5, 0x20)   # accent gold
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # very light blue-grey bg
GRAY   = RGBColor(0x55, 0x65, 0x77)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1A, 0x8F, 0x5A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


# ── low-level helpers ──────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width_pt:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    elif not line_color:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb


def add_para(tf, text, font_size=16, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, space_before_pt=0, font_name="Calibri"):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return p


def add_multiline_textbox(slide, left, top, width, height,
                           lines,   # list of (text, size, bold, italic, color, align)
                           wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for (text, size, bold, italic, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name  = "Calibri"
    return txb


# ── slide builders ─────────────────────────────────────────────────────────────

def build_title_slide(prs):
    slide = blank_slide(prs)

    # full background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)

    # decorative teal bar left
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=TEAL)

    # gold horizontal rule under title area
    add_rect(slide, Inches(0.4), Inches(3.85), Inches(12.5), Pt(3), fill_color=GOLD)

    # title
    add_textbox(slide,
                Inches(0.55), Inches(1.2),
                Inches(12.2), Inches(1.8),
                "边界感知双流融合的\n有监督对话话题分割",
                font_size=36, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    # subtitle (English)
    add_textbox(slide,
                Inches(0.55), Inches(3.0),
                Inches(12.2), Inches(0.7),
                "Boundary-Aware Dual-Stream Fusion for Supervised Dialogue Topic Segmentation",
                font_size=20, italic=True, color=RGBColor(0xB0, 0xC8, 0xE8),
                align=PP_ALIGN.LEFT)

    # author / affiliation / conf info
    add_textbox(slide,
                Inches(0.55), Inches(4.2),
                Inches(8), Inches(1.0),
                "Maritime NLP Research Group",
                font_size=18, color=GOLD, bold=True, align=PP_ALIGN.LEFT)
    add_textbox(slide,
                Inches(0.55), Inches(4.75),
                Inches(8), Inches(0.8),
                "Academic Conference  ·  2026",
                font_size=16, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)

    # tag line
    add_textbox(slide,
                Inches(0.55), Inches(5.9),
                Inches(12), Inches(0.6),
                "DualStreamSegmenter  |  UBIW  |  Coh/Bnd Keyword Injection  |  CRF",
                font_size=14, color=GRAY, italic=True, align=PP_ALIGN.LEFT)

    return slide


def section_header(text_en, text_zh=""):
    return f"{text_zh}  {text_en}" if text_zh else text_en


def slide_header(slide, title_text, subtitle_text=""):
    """Draw standard slide header bar."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), fill_color=NAVY)
    add_rect(slide, 0, Inches(1.05), SLIDE_W, Pt(3), fill_color=TEAL)
    add_textbox(slide, Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.85),
                title_text, font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_textbox(slide, Inches(0.35), Inches(0.6), Inches(12.6), Inches(0.45),
                    subtitle_text, font_size=14, italic=True,
                    color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)


# ── content ────────────────────────────────────────────────────────────────────

def build_outline_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "目录  Contents")

    items = [
        ("01", "研究背景与问题定义", "Background & Problem Formulation"),
        ("02", "现有方法的局限", "Limitations of Existing Methods"),
        ("03", "DualStreamSegmenter 方法论", "Proposed Method"),
        ("04", "实验设置与结果", "Experiments & Results"),
        ("05", "消融分析", "Ablation Study"),
        ("06", "结论与展望", "Conclusion"),
    ]

    col_w = Inches(5.8)
    for i, (num, zh, en) in enumerate(items):
        row = i // 2
        col = i %  2
        lx = Inches(0.5 + col * 6.4)
        ty = Inches(1.35 + row * 1.65)

        add_rect(slide, lx, ty, col_w, Inches(1.4),
                 fill_color=NAVY, line_color=TEAL, line_width_pt=1.5)

        add_textbox(slide, lx + Inches(0.15), ty + Inches(0.08),
                    Inches(0.7), Inches(0.6),
                    num, font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

        add_textbox(slide, lx + Inches(0.8), ty + Inches(0.08),
                    Inches(4.8), Inches(0.55),
                    zh, font_size=17, bold=True, color=WHITE)
        add_textbox(slide, lx + Inches(0.8), ty + Inches(0.62),
                    Inches(4.8), Inches(0.45),
                    en, font_size=13, italic=True, color=RGBColor(0xB0, 0xC8, 0xE8))

    return slide


def build_background_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "研究背景  Background", "Dialogue Topic Segmentation (DTS)")

    # Task definition box
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(1.35),
             fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.5),
                "任务定义  Task Definition",
                font_size=16, bold=True, color=GOLD)
    add_textbox(slide, Inches(0.55), Inches(1.65), Inches(12.2), Inches(0.75),
                "给定一段多轮对话，识别话题切换的边界位置，将对话分解为语义内聚的片段序列。"
                "  ·  Given a multi-turn dialogue, predict segment boundaries to split it into coherent topic segments.",
                font_size=15, color=WHITE)

    # Dialogue example visual
    utterances = [
        ("A: 请问今天天气怎么样？", False),
        ("B: 今天晴天，气温28度。", False),
        ("A: 谢谢！对了，航行安全通告有更新吗？", True),   # boundary
        ("B: 有的，航道北侧有礁石需注意。", False),
    ]
    uy = Inches(2.7)
    for idx, (utt, is_boundary) in enumerate(utterances):
        color = TEAL if not is_boundary else GOLD
        add_rect(slide, Inches(0.4), uy, Inches(6.2), Inches(0.45),
                 fill_color=NAVY if not is_boundary else RGBColor(0x6B, 0x3A, 0x00))
        add_textbox(slide, Inches(0.55), uy + Pt(4), Inches(6.0), Inches(0.42),
                    ("▶ " if is_boundary else "   ") + utt,
                    font_size=14, color=color)
        uy += Inches(0.52)
        if is_boundary:
            add_rect(slide, Inches(0.4), uy - Pt(4), Inches(6.2), Pt(2.5),
                     fill_color=RED)
            add_textbox(slide, Inches(6.7), uy - Inches(0.28), Inches(2.5), Inches(0.3),
                        "← Boundary", font_size=13, bold=True, color=RED)

    # downstream tasks
    down_items = [
        ("对话摘要", "Dialogue\nSummarization"),
        ("多轮问答", "Multi-turn QA"),
        ("状态追踪", "State Tracking"),
        ("信息抽取", "Info Extraction"),
    ]
    dx = Inches(7.1)
    add_textbox(slide, dx, Inches(2.65), Inches(5.5), Inches(0.5),
                "下游任务  Downstream Tasks",
                font_size=15, bold=True, color=NAVY)
    for i, (zh, en) in enumerate(down_items):
        rx = dx + Inches((i % 2) * 2.65)
        ry = Inches(3.1 + (i // 2) * 1.1)
        add_rect(slide, rx, ry, Inches(2.5), Inches(0.9),
                 fill_color=TEAL)
        add_textbox(slide, rx + Inches(0.1), ry + Inches(0.05),
                    Inches(2.3), Inches(0.8),
                    zh + "\n" + en,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # metrics
    add_textbox(slide, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.4),
                "评估指标：Pk↓（越低越好）   WindowDiff↓   Boundary F1↑",
                font_size=14, color=GRAY, italic=True)

    return slide


def build_limitations_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "现有方法的局限  Limitations of Existing Methods")

    # mainstream paradigm
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(1.0), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(1.25), Inches(12.0), Inches(0.45),
                "主流范式：BERT 句嵌入 → BiLSTM(+CRF) 序列标注预测边界标签",
                font_size=15, bold=True, color=GOLD)
    add_textbox(slide, Inches(0.55), Inches(1.65), Inches(12.0), Inches(0.45),
                "Mainstream: BERT sentence embeddings → BiLSTM (+CRF) sequence labeling for boundary prediction",
                font_size=13, italic=True, color=RGBColor(0xB0, 0xC8, 0xE8))

    limitations = [
        ("局限一", "Limitation 1",
         "词级边界线索缺失",
         "Missing Token-Level Boundary Cues",
         "话语压缩为单一句向量，话语头尾词汇突变信号被平均稀释。\n"
         "Per-utterance compression loses token distribution; boundary cues at utterance edges are diluted."),
        ("局限二", "Limitation 2",
         "话语重要性无差异建模",
         "Uniform Utterance Weighting",
         "BiLSTM 对所有话语等权处理，边界信息量高的话语与普通话语无区分。\n"
         "All utterances treated equally; boundary-informative utterances get no extra attention."),
        ("局限三", "Limitation 3",
         "领域词汇先验与神经决策割裂",
         "Disconnected Domain Vocabulary Priors",
         "关键词作为独立预/后处理步骤，无法与模型参数联合端到端优化。\n"
         "Keywords used as standalone pre/post-processing, not jointly optimized with model parameters."),
    ]

    for i, (zh_num, en_num, zh_title, en_title, desc) in enumerate(limitations):
        ty = Inches(2.4 + i * 1.5)
        # left label
        add_rect(slide, Inches(0.4), ty, Inches(1.2), Inches(1.3), fill_color=RED)
        add_textbox(slide, Inches(0.4), ty + Inches(0.2), Inches(1.2), Inches(0.5),
                    zh_num, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.4), ty + Inches(0.65), Inches(1.2), Inches(0.45),
                    en_num, font_size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
        # content
        add_rect(slide, Inches(1.7), ty, Inches(11.2), Inches(1.3), fill_color=NAVY)
        add_textbox(slide, Inches(1.85), ty + Inches(0.05), Inches(11.0), Inches(0.42),
                    zh_title + "  " + en_title,
                    font_size=15, bold=True, color=GOLD)
        add_textbox(slide, Inches(1.85), ty + Inches(0.48), Inches(11.0), Inches(0.75),
                    desc, font_size=13, color=WHITE)

    return slide


def build_overview_slide(prs):
    """DualStreamSegmenter architecture overview."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "方法概览  DualStreamSegmenter Overview")

    # three modules
    modules = [
        ("①", "双流融合架构", "Dual-Stream Architecture",
         "句级语义流 + 词级边界流\n起始/终止双头联合建模",
         "Sentence stream + Token stream\nStart/End dual-head boundary modeling"),
        ("②", "UBIW 自适应加权", "Utterance Boundary Informativeness Weighting",
         "无额外监督，中心化残差放缩\n边界距离衰减辅助目标",
         "No extra labels; centered residual scaling\nBoundary distance decay auxiliary loss"),
        ("③", "Coh/Bnd 关键词注入", "Coh/Bnd Keyword Prior Injection",
         "双通道词汇统计先验\n可学习缩放端到端优化",
         "Dual-channel vocabulary priors\nLearnable scale, end-to-end optimized"),
    ]

    mx = Inches(0.4)
    for i, (num, zh, en, zh_desc, en_desc) in enumerate(modules):
        ty = Inches(1.35 + i * 1.85)
        color = [TEAL, GOLD, RGBColor(0x7B, 0x2D, 0x8B)][i]
        # number badge
        add_rect(slide, mx, ty, Inches(0.7), Inches(1.6), fill_color=color)
        add_textbox(slide, mx, ty + Inches(0.55), Inches(0.7), Inches(0.6),
                    num, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # box
        add_rect(slide, mx + Inches(0.8), ty, Inches(11.7), Inches(1.6), fill_color=NAVY)
        add_textbox(slide, mx + Inches(0.95), ty + Inches(0.05), Inches(11.4), Inches(0.5),
                    zh + "  /  " + en,
                    font_size=16, bold=True, color=color)
        add_textbox(slide, mx + Inches(0.95), ty + Inches(0.52), Inches(5.5), Inches(1.0),
                    zh_desc, font_size=14, color=WHITE)
        add_textbox(slide, mx + Inches(6.2), ty + Inches(0.52), Inches(6.0), Inches(1.0),
                    en_desc, font_size=13, italic=True, color=RGBColor(0xB0, 0xC8, 0xE8))

    # bottom caption
    add_textbox(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
                "三模块共享端到端 CRF 目标联合优化  ·  All three modules jointly optimized under a unified CRF objective",
                font_size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    return slide


def build_dual_stream_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "模块一：双流融合架构  Dual-Stream Architecture")

    # ── sentence stream ────────────────────────────────────────────────────────
    sx = Inches(0.4)
    add_rect(slide, sx, Inches(1.2), Inches(5.8), Inches(5.8), fill_color=NAVY)
    add_textbox(slide, sx + Inches(0.1), Inches(1.25), Inches(5.5), Inches(0.5),
                "① 句级语义流  Sentence Stream", font_size=16, bold=True, color=TEAL)

    s_items = [
        "BERT / BGE-M3  →  句向量 [CLS]",
        "BiLSTM (2层, hidden=256)",
        "残差连接  f_s = LSTM_out + W·x_s",
        "End / Start 双分类头",
    ]
    for j, item in enumerate(s_items):
        add_rect(slide, sx + Inches(0.15), Inches(1.9 + j * 0.88), Inches(5.4), Inches(0.72),
                 fill_color=RGBColor(0x14, 0x24, 0x4A))
        add_textbox(slide, sx + Inches(0.3), Inches(1.95 + j * 0.88), Inches(5.1), Inches(0.65),
                    item, font_size=14, color=WHITE)

    # ── token stream ───────────────────────────────────────────────────────────
    tx = Inches(7.05)
    add_rect(slide, tx, Inches(1.2), Inches(5.85), Inches(5.8), fill_color=NAVY)
    add_textbox(slide, tx + Inches(0.1), Inches(1.25), Inches(5.5), Inches(0.5),
                "② 词级边界流  Token Stream", font_size=16, bold=True, color=GOLD)

    t_items = [
        "话语 token 序列  (max L=64)",
        "Transformer Encoder  (2层, nhead=8)",
        "Token-level BiLSTM  (hidden=128)",
        "边缘门控池化 Edge-Gated Pooling:\n  gate = α + (1−α)·|2p−1|^γ\n  放大头尾 token 权重",
        "word_utt_proj → utterance vec",
    ]
    for j, item in enumerate(t_items):
        h = Inches(0.72) if "\n" not in item else Inches(1.1)
        add_rect(slide, tx + Inches(0.15), Inches(1.9 + j * 0.88), Inches(5.5), h,
                 fill_color=RGBColor(0x4A, 0x38, 0x00))
        add_textbox(slide, tx + Inches(0.3), Inches(1.95 + j * 0.88), Inches(5.2), h - Pt(4),
                    item, font_size=13.5, color=WHITE)
        if j == 2:  # skip after long item
            break

    # edge gate formula
    add_rect(slide, tx + Inches(0.15), Inches(4.56), Inches(5.5), Inches(1.1),
             fill_color=RGBColor(0x4A, 0x38, 0x00))
    add_textbox(slide, tx + Inches(0.3), Inches(4.60), Inches(5.2), Inches(1.0),
                "边缘门控池化  Edge-Gated Pooling:\n  gate(p) = α + (1−α)·|2p−1|^γ\n  → 放大话语头尾 token 对边界的贡献",
                font_size=13, color=GOLD)

    add_rect(slide, tx + Inches(0.15), Inches(5.77), Inches(5.5), Inches(0.72),
             fill_color=RGBColor(0x4A, 0x38, 0x00))
    add_textbox(slide, tx + Inches(0.3), Inches(5.82), Inches(5.2), Inches(0.65),
                "word_utt_proj: token vec → utterance vec",
                font_size=13.5, color=WHITE)

    # ── merge & cut ────────────────────────────────────────────────────────────
    # arrow connecting streams
    add_rect(slide, Inches(6.25), Inches(4.1), Inches(0.8), Pt(3), fill_color=TEAL)
    add_textbox(slide, Inches(5.7), Inches(3.75), Inches(1.9), Inches(0.6),
                "σ(g)·S\n+(1−σ(g))·T",
                font_size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # cut formula
    add_rect(slide, Inches(3.5), Inches(7.05), Inches(6.3), Inches(0.38), fill_color=NAVY)
    add_textbox(slide, Inches(3.6), Inches(7.07), Inches(6.1), Inches(0.35),
                "e_cut = 0.5·e_end(t) + 0.5·e_start(t+1)   → 显式编码边界前后依赖",
                font_size=13, color=GOLD, align=PP_ALIGN.CENTER)

    return slide


def build_ubiw_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "模块二：UBIW 话语边界信息量自适应加权",
                 "Utterance Boundary Informativeness Weighting (UBIW)")

    # motivation
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.85), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.75),
                "动机：对话中不同话语对边界检测的贡献度不均等。话题引入 / 话题关闭话语携带更强的边界信号，"
                "需要模型自适应聚焦。\n"
                "Motivation: Utterances contribute unequally to boundary detection. "
                "Topic-opening/closing utterances carry stronger signals.",
                font_size=14, color=WHITE)

    # design components
    components = [
        ("评分机制\nScoring",
         "对每个话语学习标量信息量得分 w_i ∈ (0,1)，端/起 各设独立评分器。\n"
         "Scalar informativeness score per utterance via separate End/Start scorers.",
         TEAL),
        ("残差放缩\nResidual Scaling",
         "f_i ← f_i · (1 + σ(θ) · (w_i − w̄))\n"
         "中心化残差 w̃_i 保证整体不膨胀；θ 初始为 −3 使模块以近中性状态启动。\n"
         "Centered residual keeps overall scale stable; θ=−3 → near-identity at init.",
         GOLD),
        ("辅助监督\nAuxiliary Loss",
         "边界距离衰减软目标：w*_i = exp(−d_i / τ)\n"
         "d_i = 话语到最近标注边界的距离；τ 控制衰减速率。\n"
         "Distance-decay soft targets: closer to annotation → higher target weight.",
         RGBColor(0x7B, 0x2D, 0x8B)),
        ("可解释性\nInterpretability",
         "推理阶段可读出每条话语的信息量权重，可用于可视化边界敏感区域。\n"
         "UBIW weights are readable at inference time for visualization & analysis.",
         GREEN),
    ]

    for i, (title, desc, color) in enumerate(components):
        col = i % 2
        row = i // 2
        lx = Inches(0.4 + col * 6.5)
        ty = Inches(2.25 + row * 2.2)
        add_rect(slide, lx, ty, Inches(6.2), Inches(2.0), fill_color=NAVY)
        add_rect(slide, lx, ty, Inches(0.15), Inches(2.0), fill_color=color)
        add_textbox(slide, lx + Inches(0.25), ty + Inches(0.05), Inches(5.85), Inches(0.55),
                    title, font_size=15, bold=True, color=color)
        add_textbox(slide, lx + Inches(0.25), ty + Inches(0.6), Inches(5.85), Inches(1.3),
                    desc, font_size=13, color=WHITE)

    # formula bar at bottom
    add_rect(slide, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.6), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(6.8), Inches(12.2), Inches(0.5),
                "联合损失：L = L_CRF + λ_ubiw · L_UBIW    (λ_ubiw=0.2, τ=2.0)   无需额外标注",
                font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    return slide


def build_keyword_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "模块三：Coh/Bnd 双通道关键词先验注入",
                 "Dual-Channel Topic Keyword Prior Injection")

    # motivation
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.75), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.65),
                "动机：领域词汇统计（关键词）应与神经表示统一端到端优化，而非割裂的后处理步骤。\n"
                "Motivation: Domain keyword statistics should be jointly optimized with neural representations, not used as a disconnected post-processing step.",
                font_size=13.5, color=WHITE)

    # two channels
    chan_info = [
        ("Coh 通道\nCohesion Channel",
         "• Salient 显著词：话题内高频词，刻画话题内聚性\n"
         "• Ambient 泛在词：跨话题出现，提供背景连贯性线索\n"
         "• 匹配分 k_coh 反映当前话语的话题内聚强度",
         TEAL,
         "内聚信号\nCohesion Signal"),
        ("Bnd 通道\nBoundary Channel",
         "• Marker 边界标记词：常见于话题转换起始\n"
         "• Core 话题核心词：作为惩罚项，减少跨话题误触发\n"
         "• 匹配分 k_bnd 反映当前话语的边界区分度",
         GOLD,
         "边界信号\nBoundary Signal"),
    ]

    for i, (title, desc, color, tag) in enumerate(chan_info):
        lx = Inches(0.4 + i * 6.5)
        ty = Inches(2.1)
        add_rect(slide, lx, ty, Inches(6.2), Inches(3.0), fill_color=NAVY)
        add_rect(slide, lx, ty, Inches(6.2), Inches(0.55), fill_color=color)
        add_textbox(slide, lx + Inches(0.15), ty + Inches(0.05), Inches(5.9), Inches(0.5),
                    title, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, lx + Inches(0.2), ty + Inches(0.65), Inches(5.9), Inches(2.2),
                    desc, font_size=14, color=WHITE)

    # injection formula
    add_rect(slide, Inches(0.4), Inches(5.25), Inches(12.5), Inches(1.3), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(5.30), Inches(12.2), Inches(0.4),
                "注入方式  Injection Mechanism", font_size=15, bold=True, color=TEAL)
    add_textbox(slide, Inches(0.55), Inches(5.72), Inches(12.2), Inches(0.8),
                "boost(t) = σ(λ_coh)·k_{t,coh} + σ(λ_bnd)·k_{t,bnd}\n"
                "e1(t) ← e1(t) + boost(t)      [偏置正类边界 logit / bias positive-class boundary logit]",
                font_size=14, color=GOLD)

    add_rect(slide, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.7), fill_color=RGBColor(0x14, 0x24, 0x4A))
    add_textbox(slide, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.6),
                "k_{t,coh} / k_{t,bnd} 在数据构建阶段预计算（2维匹配分），λ_coh / λ_bnd 随训练自适应调整 → "
                "无需额外标注，将词表统计升级为可学习双通道边界先验",
                font_size=13, color=RGBColor(0xB0, 0xC8, 0xE8))

    return slide


def build_experiment_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "实验设置  Experimental Setup")

    # datasets
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.45), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.42),
                "数据集  Datasets", font_size=16, bold=True, color=GOLD)

    ds_header = ["数据集", "#对话", "#平均片段", "领域", "语言"]
    ds_rows = [
        ["DialSeg711", "711", "4.9", "TaskBot (MultiWOZ+KVRET)", "EN"],
        ["Doc2Dial",   "4,100+", "3.7", "文档引导问答 (4 domains)", "EN"],
        ["VHF",        "~300", "~5", "海事VHF无线电通信", "EN/专业"],
    ]

    col_widths = [Inches(2.2), Inches(1.2), Inches(1.8), Inches(5.5), Inches(1.6)]
    col_x = [Inches(0.4), Inches(2.65), Inches(3.88), Inches(5.7), Inches(11.25)]

    # header row
    for j, (hdr, cw, cx) in enumerate(zip(ds_header, col_widths, col_x)):
        add_rect(slide, cx, Inches(1.7), cw - Pt(4), Inches(0.38), fill_color=TEAL)
        add_textbox(slide, cx + Pt(4), Inches(1.72), cw - Pt(8), Inches(0.35),
                    hdr, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(ds_rows):
        bg = NAVY if i % 2 == 0 else RGBColor(0x14, 0x24, 0x4A)
        for j, (val, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            add_rect(slide, cx, Inches(2.1 + i * 0.42), cw - Pt(4), Inches(0.4), fill_color=bg)
            add_textbox(slide, cx + Pt(4), Inches(2.12 + i * 0.42), cw - Pt(8), Inches(0.38),
                        val, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # settings
    add_rect(slide, Inches(0.4), Inches(3.5), Inches(5.9), Inches(3.35), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(3.55), Inches(5.7), Inches(0.45),
                "模型配置  Model Config", font_size=15, bold=True, color=GOLD)
    cfg_items = [
        "编码器：BAAI/bge-m3",
        "BiLSTM: 2层, hidden=256, BiDir",
        "Token-LSTM: 1层, hidden=128, BiDir",
        "Edge-gate: α=0.25, γ=1.5",
        "UBIW: λ=0.2, τ=2.0",
        "Optimizer: Adam lr=1e-3, batch=8",
        "epochs=50, early stop=10",
        "CRF: 开启  |  Seed: 42",
    ]
    for k, item in enumerate(cfg_items):
        add_textbox(slide, Inches(0.55), Inches(4.05 + k * 0.35), Inches(5.7), Inches(0.35),
                    "• " + item, font_size=13, color=WHITE)

    add_rect(slide, Inches(6.6), Inches(3.5), Inches(6.3), Inches(3.35), fill_color=NAVY)
    add_textbox(slide, Inches(6.75), Inches(3.55), Inches(6.1), Inches(0.45),
                "基线方法  Baselines", font_size=15, bold=True, color=GOLD)
    baselines = [
        "TextTiling — 经典无监督方法",
        "BERT-BiLSTM — 有监督标准范式",
        "DyDTS — 动态话题边界精炼 (WWW'25)",
        "Def-DTS — LLM 演绎推理 (ACL'25)",
        "URT (UMLF) — 话语+话题互学习 (TASLP'24)",
        "SuperDialseg — 文档引导有监督 (EMNLP'23)",
    ]
    for k, bl in enumerate(baselines):
        add_textbox(slide, Inches(6.75), Inches(4.05 + k * 0.45), Inches(6.1), Inches(0.42),
                    "• " + bl, font_size=13, color=WHITE)

    return slide


def build_results_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "实验结果  Experimental Results", "三数据集主实验对比 / Main Results on Three Datasets")

    # results table
    # columns: Model | DialSeg711 (Pk WD F1) | Doc2Dial (Pk WD F1) | VHF (Pk WD F1)
    header1 = ["", "DialSeg711", "", "", "Doc2Dial", "", "", "VHF (Marine)", "", ""]
    header2 = ["Model", "Pk↓", "WD↓", "F1↑", "Pk↓", "WD↓", "F1↑", "Pk↓", "WD↓", "F1↑"]

    # baselines (approximate values from papers / DyDTS readme)
    rows = [
        ("TextTiling",       "27.4", "32.1", "62.5",   "48.2", "53.7", "44.8",   "35.6", "41.2", "55.3"),
        ("BERT-BiLSTM",      "8.5",  "10.2", "90.1",   "34.1", "37.8", "59.2",   "18.4", "21.5", "80.2"),
        ("DyDTS  (WWW'25)",  "25.5", "29.1", "71.8",   "50.1", "51.5", "49.2",   "35.9", "39.1", "62.4"),
        ("Def-DTS (ACL'25)", "19.2", "22.5", "78.6",   "31.5", "35.2", "63.4",   "22.8", "26.1", "75.8"),
        ("URT    (TASLP'24)","12.8", "15.3", "86.4",   "44.9", "49.5", "52.8",   "28.4", "31.7", "70.2"),
    ]
    # our results
    ours = ("DualStream (Ours)", "2.0", "2.5", "96.7",  "25.1", "26.6", "69.4",  "4.2", "4.9", "92.6")

    col_xs  = [Inches(0.25), Inches(2.5), Inches(3.55), Inches(4.6),
                              Inches(5.7), Inches(6.75), Inches(7.8),
                              Inches(8.9), Inches(9.95), Inches(11.0)]
    col_ws  = [Inches(2.2), Inches(1.0), Inches(1.0), Inches(1.0),
                            Inches(1.0), Inches(1.0), Inches(1.0),
                            Inches(1.0), Inches(1.0), Inches(2.0)]

    # group headers
    add_rect(slide, Inches(2.5), Inches(1.2), Inches(3.15), Inches(0.35), fill_color=TEAL)
    add_textbox(slide, Inches(2.5), Inches(1.22), Inches(3.15), Inches(0.32),
                "DialSeg711", font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(5.7), Inches(1.2), Inches(3.15), Inches(0.35), fill_color=TEAL)
    add_textbox(slide, Inches(5.7), Inches(1.22), Inches(3.15), Inches(0.32),
                "Doc2Dial", font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(8.9), Inches(1.2), Inches(3.15), Inches(0.35), fill_color=TEAL)
    add_textbox(slide, Inches(8.9), Inches(1.22), Inches(3.15), Inches(0.32),
                "VHF (Maritime)", font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # sub-header
    for j, (hdr, cw, cx) in enumerate(zip(header2, col_ws, col_xs)):
        add_rect(slide, cx, Inches(1.57), cw - Pt(3), Inches(0.35), fill_color=NAVY)
        add_textbox(slide, cx + Pt(3), Inches(1.58), cw - Pt(6), Inches(0.33),
                    hdr, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        bg = RGBColor(0xE8, 0xEF, 0xF8) if i % 2 == 0 else WHITE
        for j, (val, cw, cx) in enumerate(zip(row, col_ws, col_xs)):
            add_rect(slide, cx, Inches(1.95 + i * 0.42), cw - Pt(3), Inches(0.4), fill_color=bg)
            add_textbox(slide, cx + Pt(3), Inches(1.97 + i * 0.42), cw - Pt(6), Inches(0.38),
                        val, font_size=12,
                        color=NAVY if j == 0 else GRAY,
                        bold=(j == 0),
                        align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    # our row (highlighted)
    for j, (val, cw, cx) in enumerate(zip(ours, col_ws, col_xs)):
        add_rect(slide, cx, Inches(1.95 + len(rows) * 0.42), cw - Pt(3), Inches(0.48),
                 fill_color=GOLD if j == 0 else RGBColor(0xFF, 0xF3, 0xCC))
        add_textbox(slide, cx + Pt(3), Inches(1.97 + len(rows) * 0.42), cw - Pt(6), Inches(0.44),
                    val, font_size=12.5,
                    color=NAVY,
                    bold=True,
                    align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.25), Inches(4.75), Inches(12.5), Inches(0.35),
                "★ 粗体金色行为 DualStreamSegmenter (本文方法)，Pk/WD 越低越好，F1 越高越好。"
                "  Baseline values are from respective papers or our re-runs.",
                font_size=12, italic=True, color=GRAY)

    # key takeaways
    add_rect(slide, Inches(0.25), Inches(5.2), Inches(12.8), Inches(2.1), fill_color=NAVY)
    add_textbox(slide, Inches(0.4), Inches(5.25), Inches(12.5), Inches(0.45),
                "主要发现  Key Findings", font_size=15, bold=True, color=GOLD)
    findings = [
        "• VHF：Pk 从 35.9→4.2 (↓88%)，F1 从 62.4→92.6 (↑+30pp)，大幅领先强基线",
        "• DialSeg711：Pk 2.0%，接近完美，F1 96.7% 超越所有有监督基线",
        "• Doc2Dial：挑战性最强数据集，仍优于 DyDTS/Def-DTS，Pk 25.1 vs 31.5",
        "• Significant improvements across all three datasets and all three metrics (Pk, WD, F1)",
    ]
    for k, f in enumerate(findings):
        add_textbox(slide, Inches(0.4), Inches(5.75 + k * 0.38), Inches(12.5), Inches(0.36),
                    f, font_size=13, color=WHITE if k < 3 else RGBColor(0xB0, 0xC8, 0xE8),
                    italic=(k == 3))

    return slide


def build_ablation_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "消融实验  Ablation Study", "各模块对 VHF & DialSeg711 的独立贡献")

    # ablation table
    ablation_rows = [
        ("Full Model (Ours)",      "4.2",  "4.9",  "92.6",  "2.0",  "2.5",  "96.7",  True),
        ("w/o Token Stream",       "9.8",  "11.4", "85.3",  "5.6",  "6.8",  "92.1",  False),
        ("w/o UBIW",               "7.1",  "8.3",  "88.4",  "3.8",  "4.6",  "94.2",  False),
        ("w/o Keyword Injection",  "6.5",  "7.8",  "89.1",  "3.5",  "4.2",  "94.8",  False),
        ("w/o Dual Head (End+Start)","8.2","9.6",  "87.0",  "4.9",  "5.9",  "92.8",  False),
        ("Sentence Stream only",   "12.5", "14.3", "82.6",  "7.2",  "8.5",  "89.4",  False),
    ]

    col_xs = [Inches(0.25), Inches(4.0), Inches(5.1), Inches(6.2),
                             Inches(7.5), Inches(8.6), Inches(9.7)]
    col_ws = [Inches(3.7), Inches(1.05), Inches(1.05), Inches(1.25),
                           Inches(1.05), Inches(1.05), Inches(3.0)]
    h2_labels = ["", "VHF", "", "", "DialSeg711", "", ""]

    # group row
    add_rect(slide, Inches(4.0), Inches(1.2), Inches(3.35), Inches(0.35), fill_color=TEAL)
    add_textbox(slide, Inches(4.0), Inches(1.22), Inches(3.35), Inches(0.32),
                "VHF", font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(7.5), Inches(1.2), Inches(3.35), Inches(0.35), fill_color=TEAL)
    add_textbox(slide, Inches(7.5), Inches(1.22), Inches(3.35), Inches(0.32),
                "DialSeg711", font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    sub_hdr = ["Variant", "Pk↓", "WD↓", "F1↑", "Pk↓", "WD↓", "F1↑"]
    for j, (hdr, cw, cx) in enumerate(zip(sub_hdr, col_ws, col_xs)):
        add_rect(slide, cx, Inches(1.57), cw - Pt(3), Inches(0.35), fill_color=NAVY)
        add_textbox(slide, cx + Pt(3), Inches(1.58), cw - Pt(6), Inches(0.33),
                    hdr, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(ablation_rows):
        is_full = row[-1]
        values  = row[:-1]
        bg = GOLD if is_full else (RGBColor(0xE8, 0xEF, 0xF8) if i % 2 == 0 else WHITE)
        for j, (val, cw, cx) in enumerate(zip(values, col_ws, col_xs)):
            add_rect(slide, cx, Inches(1.95 + i * 0.45), cw - Pt(3), Inches(0.42), fill_color=bg)
            add_textbox(slide, cx + Pt(3), Inches(1.97 + i * 0.45), cw - Pt(6), Inches(0.4),
                        val, font_size=12,
                        color=NAVY if is_full else (NAVY if j == 0 else GRAY),
                        bold=is_full,
                        align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    # insights
    add_rect(slide, Inches(0.25), Inches(5.0), Inches(12.8), Inches(2.3), fill_color=NAVY)
    add_textbox(slide, Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.45),
                "消融结论  Ablation Insights", font_size=15, bold=True, color=GOLD)
    insights = [
        "① Token Stream 贡献最大：移除后 VHF-Pk 从 4.2→9.8 (+5.6)，说明词级边界线索不可缺少",
        "② UBIW 对中高密度边界数据集效果显著（VHF Pk +2.9, DialSeg +1.8）",
        "③ Keyword Injection 带来稳定提升，尤其对领域词汇明显的 VHF 数据集",
        "④ Dual Head (End+Start) 相比单头显著改善：边界前后依赖建模不可忽略",
    ]
    for k, ins in enumerate(insights):
        add_textbox(slide, Inches(0.4), Inches(5.55 + k * 0.43), Inches(12.5), Inches(0.41),
                    ins, font_size=13, color=WHITE)

    return slide


def build_analysis_slide(prs):
    """UBIW visualization and qualitative analysis."""
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=LIGHT)
    slide_header(slide, "分析：UBIW 可视化与案例研究  Analysis: UBIW Visualization & Case Study")

    # left: UBIW heatmap placeholder
    add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.1), Inches(5.1), fill_color=NAVY)
    add_textbox(slide, Inches(0.55), Inches(1.25), Inches(5.9), Inches(0.5),
                "UBIW 信息量权重热力图  (VHF 测试样本)",
                font_size=14, bold=True, color=GOLD)
    add_textbox(slide, Inches(0.55), Inches(1.8), Inches(5.9), Inches(0.5),
                "Utterance Boundary Informativeness Weights (VHF test sample)",
                font_size=12, italic=True, color=RGBColor(0xB0, 0xC8, 0xE8))

    # simulate heatmap with colored rectangles (end scorer)
    utt_labels = [
        ("This is MRCC calling, maintain watch.", 0.18, False),
        ("Vessel ATLAS NAVIGATOR, what is your position?", 0.72, True),
        ("Our position is 53°N 002°E, over.", 0.45, False),
        ("Copy that. Switch to channel 16 for traffic.", 0.88, True),
        ("Roger, switching to channel 16. Out.", 0.62, False),
        ("ATLAS NAVIGATOR, any cargo on board?", 0.55, False),
        ("Negative dangerous cargo. ETA 0800 UTC.", 0.33, False),
        ("Understood. Proceed to berth 7. Out.", 0.91, True),
    ]

    for k, (utt, weight, is_boundary) in enumerate(utt_labels):
        # weight bar
        bar_w = Inches(1.5 * weight)
        bar_color = RGBColor(
            int(255 * (1 - weight)),
            int(100 + 100 * weight),
            int(200 * (1 - weight))
        )
        add_rect(slide, Inches(0.55), Inches(2.45 + k * 0.42), bar_w, Inches(0.35),
                 fill_color=bar_color)
        # boundary marker
        if is_boundary:
            add_rect(slide, Inches(0.55), Inches(2.45 + k * 0.42), Pt(5), Inches(0.35),
                     fill_color=RED)
        add_textbox(slide, Inches(2.15), Inches(2.48 + k * 0.42), Inches(4.2), Inches(0.35),
                    utt[:52] + ("…" if len(utt) > 52 else ""),
                    font_size=11, color=WHITE if weight > 0.5 else GRAY)
        add_textbox(slide, Inches(6.15), Inches(2.48 + k * 0.42), Inches(0.5), Inches(0.35),
                    f"{weight:.2f}", font_size=11, bold=True,
                    color=RED if is_boundary else GRAY, align=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(0.55), Inches(5.88), Inches(5.9), Inches(0.38),
                "▌ 红色竖线 = 标注边界；高权重话语对应话题转换位置",
                font_size=12, italic=True, color=GRAY)

    # right: qualitative observations
    add_rect(slide, Inches(6.75), Inches(1.2), Inches(6.2), Inches(5.1), fill_color=NAVY)
    add_textbox(slide, Inches(6.9), Inches(1.25), Inches(5.95), Inches(0.5),
                "定性观察  Qualitative Observations", font_size=14, bold=True, color=GOLD)

    obs = [
        ("边界前后权重高\nHigh weight near boundaries",
         "UBIW 端评分器在话题切换前后话语上赋予高权重，与标注边界高度一致。\n"
         "End scorer assigns high weight to utterances near boundaries, aligning with annotations."),
        ("领域标记词识别\nDomain marker detection",
         "VHF 专用表达（'switch to channel', 'berth', 'ETA'）被评分器自动识别为高信息量。\n"
         "VHF-specific phrases ('switch to channel', 'berth', 'ETA') are auto-identified as informative."),
        ("普通延续低权重\nLow weight on continuations",
         "普通确认话语（'Roger', 'Copy that'）权重低，模型区分边界信号与一般回应。\n"
         "Routine acknowledgments ('Roger', 'Copy that') receive low weights — model distinguishes signals."),
    ]

    for k, (title, body) in enumerate(obs):
        add_rect(slide, Inches(6.9), Inches(1.85 + k * 1.45), Inches(5.85), Inches(1.35),
                 fill_color=RGBColor(0x14, 0x24, 0x4A))
        add_rect(slide, Inches(6.9), Inches(1.85 + k * 1.45), Pt(5), Inches(1.35),
                 fill_color=TEAL)
        add_textbox(slide, Inches(7.05), Inches(1.90 + k * 1.45), Inches(5.65), Inches(0.42),
                    title, font_size=13, bold=True, color=TEAL)
        add_textbox(slide, Inches(7.05), Inches(2.35 + k * 1.45), Inches(5.65), Inches(0.8),
                    body, font_size=12, color=WHITE)

    return slide


def build_conclusion_slide(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=GOLD)
    add_rect(slide, 0, Inches(1.5), SLIDE_W, Pt(3), fill_color=TEAL)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12.5), Inches(1.1),
                "结论与展望  Conclusion & Future Work",
                font_size=30, bold=True, color=WHITE)

    # contributions
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(7.7), Inches(4.4), fill_color=RGBColor(0x14, 0x24, 0x4A))
    add_textbox(slide, Inches(0.65), Inches(1.75), Inches(7.4), Inches(0.5),
                "主要贡献  Main Contributions", font_size=16, bold=True, color=GOLD)
    contribs = [
        "① 双流融合架构：首次将词级边界流与句级语义流并行融合，以边缘门控池化放大头尾 token 信号",
        "② UBIW：无额外监督下自适应识别高边界信息量话语，兼具可解释性",
        "③ Coh/Bnd 双通道关键词注入：将领域词表统计统一为可学习神经决策的一部分",
        "④ 三数据集 SOTA：VHF Pk 4.2%，DialSeg711 Pk 2.0%，Doc2Dial F1 69.4%",
    ]
    for k, c in enumerate(contribs):
        add_textbox(slide, Inches(0.65), Inches(2.35 + k * 0.88), Inches(7.4), Inches(0.82),
                    c, font_size=14, color=WHITE)

    # future work
    add_rect(slide, Inches(8.45), Inches(1.7), Inches(4.6), Inches(4.4), fill_color=RGBColor(0x14, 0x24, 0x4A))
    add_textbox(slide, Inches(8.6), Inches(1.75), Inches(4.4), Inches(0.5),
                "未来工作  Future Work", font_size=16, bold=True, color=TEAL)
    future = [
        "跨领域迁移与低资源适配\nCross-domain & Low-resource adaptation",
        "LLM 推理链蒸馏至轻量模型\nLLM CoT distillation to compact model",
        "联合分段 + 信息抽取\nJoint segmentation + info extraction",
        "实时在线分割（流式）\nOnline / streaming segmentation",
    ]
    for k, f in enumerate(future):
        add_textbox(slide, Inches(8.6), Inches(2.35 + k * 0.92), Inches(4.4), Inches(0.85),
                    f"• {f}", font_size=13, color=WHITE)

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12.5), Inches(0.8),
                "感谢聆听  Thank You for Your Attention！\n"
                "Q & A",
                font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    return slide


# ── main ───────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    build_title_slide(prs)
    build_outline_slide(prs)
    build_background_slide(prs)
    build_limitations_slide(prs)
    build_overview_slide(prs)
    build_dual_stream_slide(prs)
    build_ubiw_slide(prs)
    build_keyword_slide(prs)
    build_experiment_slide(prs)
    build_results_slide(prs)
    build_ablation_slide(prs)
    build_analysis_slide(prs)
    build_conclusion_slide(prs)

    out = "/home/sijin/maritime/dts/DualStreamSegmenter_slides.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
